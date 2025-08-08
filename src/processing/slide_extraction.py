import os
import io
import base64
import uuid
import subprocess
import tempfile
from pathlib import Path
from typing import List, Optional, Dict, Any
from PIL import Image
import logging

# Importações para PowerPoint
try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    Presentation = None

# Importações para PDF
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

from ..config.settings import settings

logger = logging.getLogger(__name__)


class SlideExtractionService:
    """
    Serviço para extrair slides de apresentações PowerPoint e PDF
    Foca na conversão para imagens para análise visual com Gemini
    """
    
    def __init__(self):
        self.temp_dir = settings.temp_extraction_directory
        os.makedirs(self.temp_dir, exist_ok=True)
        
        # Configurações de qualidade de imagem
        self.image_dpi = 300 if settings.image_quality == "high" else 200 if settings.image_quality == "medium" else 150
        self.image_format = "PNG"  # PNG para melhor qualidade
        self.max_image_dimension = 1920  # Máximo para uma dimensão
    
    async def extract_slides_from_file(self, file_path: str) -> Dict[str, Any]:
        """
        Extrai slides de um arquivo de apresentação
        Foca na conversão para imagens para análise visual
        
        Returns:
            Dict com:
            - slides_data: Lista de dados dos slides (simplificada)
            - metadata: Metadados da apresentação
        """
        file_extension = Path(file_path).suffix.lower()
        
        if file_extension in ['.pptx', '.ppt']:
            return await self._extract_from_powerpoint(file_path)
        elif file_extension == '.pdf':
            return await self._extract_from_pdf(file_path)
        else:
            raise ValueError(f"Formato de arquivo não suportado: {file_extension}")
    
    async def _extract_from_powerpoint(self, file_path: str) -> Dict[str, Any]:
        """
        Extrai slides de arquivo PowerPoint
        Converte cada slide para imagem para análise visual
        """
        if Presentation is None:
            raise ImportError("python-pptx não está instalado. Execute: pip install python-pptx")
        
        try:
            # Carrega apresentação apenas para metadados básicos
            prs = Presentation(file_path)
            
            # Extrai metadados
            metadata = {
                "total_slides": len(prs.slides),
                "title": getattr(prs.core_properties, 'title', None) or '',
                "author": getattr(prs.core_properties, 'author', None) or '',
                "created": getattr(prs.core_properties, 'created', None),
                "modified": getattr(prs.core_properties, 'modified', None),
                "source_filename": Path(file_path).name
            }
            
            # Converte PowerPoint para imagens usando LibreOffice
            slide_images = await self._convert_pptx_to_images(file_path)
            
            # Cria dados simplificados dos slides
            slides_data = []
            for slide_idx in range(len(prs.slides)):
                slide_number = slide_idx + 1
                image_path = slide_images.get(slide_number)
                
                slide_data = {
                    "slide_number": slide_number,
                    "image_path": image_path,
                    "image_base64": None,  # Será preenchido quando necessário
                    "layout_name": None,  # Será determinado pela análise visual
                    "extracted_text": None,  # Será extraído pelo Gemini
                    "slide_title": None,  # Será determinado pelo Gemini
                }
                
                # Converte imagem para base64 imediatamente
                if image_path and os.path.exists(image_path):
                    slide_data["image_base64"] = await self.convert_image_to_base64(image_path)
                
                slides_data.append(slide_data)
                logger.info(f"Slide {slide_number}/{len(prs.slides)} processado e convertido para imagem")
            
            return {
                "slides_data": slides_data,
                "metadata": metadata
            }
            
        except Exception as e:
            logger.error(f"Erro ao extrair slides do PowerPoint: {e}")
            raise
    
    async def _convert_pptx_to_images(self, file_path: str) -> Dict[int, str]:
        """
        Converte arquivo PowerPoint para imagens via PDF
        Estratégia: PPTX -> PDF (LibreOffice) -> PNG por página (PyMuPDF)
        """
        slide_images = {}
        
        try:
            # Cria diretório temporário para conversão
            temp_output_dir = os.path.join(self.temp_dir, f"pptx_export_{uuid.uuid4().hex[:8]}")
            os.makedirs(temp_output_dir, exist_ok=True)
            
            # Primeira etapa: Converte PPTX para PDF usando LibreOffice
            pdf_path = await self._convert_pptx_to_pdf(file_path, temp_output_dir)
            
            if pdf_path and os.path.exists(pdf_path):
                logger.info(f"PPTX convertido para PDF: {pdf_path}")
                
                # Segunda etapa: Extrai páginas do PDF como imagens usando PyMuPDF
                if fitz is not None:
                    slide_images = await self._extract_images_from_pdf(pdf_path, file_path)
                    logger.info(f"Extraídas {len(slide_images)} imagens do PDF convertido")
                else:
                    logger.warning("PyMuPDF não disponível para extrair imagens do PDF")
                    slide_images = await self._fallback_pptx_to_images(file_path)
                
                # Remove PDF temporário
                try:
                    os.unlink(pdf_path)
                except:
                    pass
            else:
                logger.warning("Falha na conversão PPTX->PDF, tentando método alternativo")
                slide_images = await self._fallback_pptx_to_images(file_path)
            
            # Remove diretório temporário
            import shutil
            shutil.rmtree(temp_output_dir, ignore_errors=True)
            
        except Exception as e:
            logger.error(f"Erro na conversão PPTX->PDF->PNG: {e}")
            slide_images = await self._fallback_pptx_to_images(file_path)
        
        return slide_images
    
    async def _convert_pptx_to_pdf(self, file_path: str, output_dir: str) -> Optional[str]:
        """
        Converte arquivo PPTX para PDF usando LibreOffice
        """
        try:
            # Comando LibreOffice para conversão para PDF
            cmd = [
                "libreoffice",
                "--headless",
                "--convert-to", "pdf",
                "--outdir", output_dir,
                file_path
            ]
            
            # Executa conversão
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            
            if result.returncode == 0:
                # Determina nome do arquivo PDF gerado
                base_name = Path(file_path).stem
                pdf_path = os.path.join(output_dir, f"{base_name}.pdf")
                
                if os.path.exists(pdf_path):
                    return pdf_path
                else:
                    logger.warning(f"PDF esperado não encontrado: {pdf_path}")
                    return None
            else:
                logger.warning(f"LibreOffice falhou na conversão para PDF: {result.stderr}")
                return None
                
        except subprocess.TimeoutExpired:
            logger.error("Timeout na conversão PPTX->PDF com LibreOffice")
            return None
        except Exception as e:
            logger.error(f"Erro na conversão PPTX->PDF: {e}")
            return None
    
    async def _extract_images_from_pdf(self, pdf_path: str, original_pptx_path: str) -> Dict[int, str]:
        """
        Extrai imagens de cada página do PDF convertido
        """
        slide_images = {}
        
        try:
            doc = fitz.open(pdf_path)
            
            for page_idx in range(len(doc)):
                page_number = page_idx + 1
                page = doc[page_idx]
                
                # Usa o método existente para salvar página como imagem
                image_path = await self._save_page_as_image_pdf(page, page_number, original_pptx_path)
                
                if image_path and os.path.exists(image_path):
                    slide_images[page_number] = image_path
                    logger.debug(f"Página {page_number} extraída como imagem: {image_path}")
                else:
                    logger.warning(f"Falha ao extrair imagem da página {page_number}")
            
            doc.close()
            return slide_images
            
        except Exception as e:
            logger.error(f"Erro ao extrair imagens do PDF: {e}")
            return {}
    
    async def _fallback_pptx_to_images(self, file_path: str) -> Dict[int, str]:
        """
        Método de fallback para converter PPTX em imagens
        Melhorado para lidar com diferentes formatos de imagem
        """
        logger.info("Usando método de fallback para conversão de slides")
        
        try:
            prs = Presentation(file_path)
            slide_images = {}
            
            # Cria imagens placeholder ou tenta extrair imagens existentes
            for slide_idx, slide in enumerate(prs.slides, 1):
                extracted_image = False
                
                # Busca por imagens no slide
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        try:
                            # Extrai imagem existente do slide
                            image = shape.image
                            image_bytes = image.blob
                            
                            # Detecta formato da imagem
                            image_format = self._detect_image_format(image_bytes)
                            
                            if image_format in ['JPEG', 'PNG', 'BMP', 'GIF']:
                                # Salva imagem extraída
                                image_path = os.path.join(
                                    self.temp_dir,
                                    f"{Path(file_path).stem}_slide_{slide_idx}_extracted_{uuid.uuid4().hex[:8]}.{image_format.lower()}"
                                )
                                
                                with open(image_path, 'wb') as f:
                                    f.write(image_bytes)
                                
                                # Converte para PNG se necessário e valida
                                final_path = await self._ensure_compatible_format(image_path)
                                if final_path:
                                    slide_images[slide_idx] = final_path
                                    extracted_image = True
                                    break
                            else:
                                logger.debug(f"Formato de imagem não suportado no slide {slide_idx}: {image_format}")
                                
                        except Exception as e:
                            logger.debug(f"Erro ao extrair imagem do slide {slide_idx}: {e}")
                            continue
                
                # Se não encontrou imagem válida, cria placeholder
                if not extracted_image:
                    placeholder_path = await self._create_placeholder_image(slide_idx, file_path)
                    if placeholder_path:
                        slide_images[slide_idx] = placeholder_path
            
            return slide_images
            
        except Exception as e:
            logger.error(f"Erro no método de fallback: {e}")
            return {}
    
    def _detect_image_format(self, image_bytes: bytes) -> str:
        """
        Detecta o formato de uma imagem pelos bytes iniciais
        """
        if image_bytes.startswith(b'\xff\xd8\xff'):
            return 'JPEG'
        elif image_bytes.startswith(b'\x89PNG'):
            return 'PNG'
        elif image_bytes.startswith(b'BM'):
            return 'BMP'
        elif image_bytes.startswith(b'GIF'):
            return 'GIF'
        elif image_bytes.startswith(b'\xd7\xcd\xc6\x9a'):
            return 'WMF'  # Windows Metafile
        elif image_bytes.startswith(b'\x01\x00\x09\x00'):
            return 'EMF'  # Enhanced Metafile
        else:
            return 'UNKNOWN'
    
    async def _ensure_compatible_format(self, image_path: str) -> Optional[str]:
        """
        Garante que a imagem está em formato compatível (PNG/JPEG)
        """
        try:
            with Image.open(image_path) as img:
                # Se já é um formato compatível, retorna o path original
                if img.format in ['PNG', 'JPEG']:
                    return image_path
                
                # Converte para PNG
                new_path = image_path.rsplit('.', 1)[0] + '_converted.png'
                
                # Trata transparência adequadamente
                if img.mode in ('RGBA', 'LA') or 'transparency' in img.info:
                    img.save(new_path, 'PNG')
                else:
                    # Converte para RGB e salva como JPEG (menor tamanho)
                    rgb_img = img.convert('RGB')
                    new_path = new_path.rsplit('.', 1)[0] + '.jpg'
                    rgb_img.save(new_path, 'JPEG', quality=90)
                
                # Remove arquivo original se conversão foi bem-sucedida
                try:
                    os.unlink(image_path)
                except:
                    pass
                
                return new_path
                
        except Exception as e:
            logger.debug(f"Erro ao converter formato de imagem: {e}")
            # Se falhar, remove arquivo problemático
            try:
                os.unlink(image_path)
            except:
                pass
            return None
    
    async def _create_placeholder_image(self, slide_number: int, file_path: str) -> Optional[str]:
        """
        Cria uma imagem placeholder para um slide
        """
        try:
            # Cria imagem simples com número do slide
            img = Image.new('RGB', (1920, 1080), color='white')
            
            # Adiciona texto simples (se PIL tiver suporte a fontes)
            try:
                from PIL import ImageDraw, ImageFont
                draw = ImageDraw.Draw(img)
                
                # Tenta usar fonte do sistema
                try:
                    font = ImageFont.truetype("/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf", 48)
                except:
                    font = ImageFont.load_default()
                
                text = f"Slide {slide_number}\n\n(Imagem não disponível)\n\nAnálise baseada em conteúdo estrutural"
                bbox = draw.textbbox((0, 0), text, font=font)
                text_width = bbox[2] - bbox[0]
                text_height = bbox[3] - bbox[1]
                
                x = (1920 - text_width) // 2
                y = (1080 - text_height) // 2
                
                draw.text((x, y), text, fill='black', font=font)
                
            except ImportError:
                # Se não conseguir adicionar texto, deixa em branco
                pass
            
            # Salva placeholder
            placeholder_path = os.path.join(
                self.temp_dir,
                f"{Path(file_path).stem}_slide_{slide_number}_placeholder_{uuid.uuid4().hex[:8]}.png"
            )
            
            img.save(placeholder_path, self.image_format)
            return placeholder_path
            
        except Exception as e:
            logger.error(f"Erro ao criar placeholder: {e}")
            return None
    
    def cleanup_temp_files(self, slide_images: List[str]):
        """
        Remove arquivos temporários de imagens
        """
        for image_path in slide_images:
            if image_path and os.path.exists(image_path):
                try:
                    os.unlink(image_path)
                except Exception as e:
                    logger.warning(f"Erro ao remover arquivo temporário {image_path}: {e}")
    
    def get_slide_image_paths(self, slides_data: List[Dict]) -> List[str]:
        """
        Extrai caminhos de imagens dos dados dos slides para limpeza
        """
        image_paths = []
        for slide in slides_data:
            if slide.get("image_path"):
                image_paths.append(slide["image_path"])
        return image_paths
    
    async def validate_image_quality(self, image_path: str) -> Dict[str, Any]:
        """
        Valida a qualidade da imagem extraída
        """
        try:
            with Image.open(image_path) as img:
                return {
                    "width": img.width,
                    "height": img.height,
                    "mode": img.mode,
                    "format": img.format,
                    "size_bytes": os.path.getsize(image_path),
                    "aspect_ratio": img.width / img.height if img.height > 0 else 0,
                    "is_valid": True
                }
        except Exception as e:
            logger.error(f"Erro ao validar imagem {image_path}: {e}")
            return {"is_valid": False, "error": str(e)}
    
    async def _extract_from_pdf(self, file_path: str) -> Dict[str, Any]:
        """
        Extrai slides de arquivo PDF
        Converte cada página para imagem para análise visual
        """
        if fitz is None:
            raise ImportError("PyMuPDF não está instalado. Execute: pip install PyMuPDF")
        
        try:
            # Abre PDF
            doc = fitz.open(file_path)
            
            # Extrai metadados
            metadata = {
                "total_slides": len(doc),
                "title": doc.metadata.get('title', ''),
                "author": doc.metadata.get('author', ''),
                "created": doc.metadata.get('creationDate', None),
                "modified": doc.metadata.get('modDate', None),
                "source_filename": Path(file_path).name
            }
            
            # Converte páginas para imagens
            slides_data = []
            
            for page_idx in range(len(doc)):
                page_number = page_idx + 1
                logger.info(f"Convertendo página {page_number}/{len(doc)} para imagem")
                
                page = doc[page_idx]
                
                # Converte página para imagem
                image_path = await self._save_page_as_image_pdf(page, page_number, file_path)
                
                slide_data = {
                    "slide_number": page_number,
                    "image_path": image_path,
                    "image_base64": None,  # Será preenchido quando necessário
                    "layout_name": "PDF Page",
                    "extracted_text": None,  # Será extraído pelo Gemini
                    "slide_title": None,  # Será determinado pelo Gemini
                }
                
                # Converte imagem para base64 imediatamente
                if image_path and os.path.exists(image_path):
                    slide_data["image_base64"] = await self.convert_image_to_base64(image_path)
                
                slides_data.append(slide_data)
            
            doc.close()
            
            return {
                "slides_data": slides_data,
                "metadata": metadata
            }
            
        except Exception as e:
            logger.error(f"Erro ao extrair slides do PDF: {e}")
            raise
    
    async def _extract_pdf_page(self, page, page_number: int, pdf_path: str) -> Dict[str, Any]:
        """
        Extrai dados de uma página específica do PDF
        """
        slide_data = {
            "slide_number": page_number,
            "title": "",
            "text_content": [],
            "layout_name": "PDF Page",
            "image_path": None,
            "shapes": []
        }
        
        # Extrai texto da página
        text_content = page.get_text()
        if text_content.strip():
            lines = [line.strip() for line in text_content.split('\n') if line.strip()]
            slide_data["text_content"] = lines
            
            # Primeira linha significativa vira título
            if lines:
                slide_data["title"] = lines[0][:100]
        
        # Extrai texto estruturado (blocos)
        blocks = page.get_text("dict")["blocks"]
        for block in blocks:
            if "lines" in block:
                block_text = ""
                for line in block["lines"]:
                    for span in line["spans"]:
                        block_text += span["text"] + " "
                
                if block_text.strip():
                    shape_data = {
                        "shape_type": "text_block",
                        "left": block["bbox"][0],
                        "top": block["bbox"][1],
                        "width": block["bbox"][2] - block["bbox"][0],
                        "height": block["bbox"][3] - block["bbox"][1],
                        "text": block_text.strip()
                    }
                    slide_data["shapes"].append(shape_data)
        
        # Salva página como imagem
        try:
            image_path = await self._save_page_as_image_pdf(page, page_number, pdf_path)
            slide_data["image_path"] = image_path
        except Exception as e:
            logger.warning(f"Não foi possível salvar imagem da página {page_number}: {e}")
        
        return slide_data
    
    async def _save_slide_as_image_pptx(self, slide, slide_number: int, presentation_path: str) -> str:
        """
        Salva slide como imagem (PowerPoint é mais complexo, precisaria de conversão)
        Por enquanto, retorna None - implementação futura
        """
        # Para PowerPoint, seria necessário usar bibliotecas como comtypes (Windows)
        # ou convert via LibreOffice headless
        # Por simplicidade, retornamos None por enquanto
        return None
    
    async def _save_page_as_image_pdf(self, page, page_number: int, pdf_path: str) -> str:
        """
        Salva página PDF como imagem com qualidade otimizada
        """
        # Calcula zoom baseado na configuração de DPI
        zoom = self.image_dpi / 72.0  # 72 DPI é o padrão do PDF
        mat = fitz.Matrix(zoom, zoom)
        
        # Renderiza página como imagem
        pix = page.get_pixmap(matrix=mat, alpha=False)
        
        # Converte para PIL Image para processamento adicional
        img_data = pix.tobytes("png")
        pil_image = Image.open(io.BytesIO(img_data))
        
        # Redimensiona se necessário para respeitar limite máximo
        if max(pil_image.size) > self.max_image_dimension:
            ratio = self.max_image_dimension / max(pil_image.size)
            new_size = tuple(int(dim * ratio) for dim in pil_image.size)
            pil_image = pil_image.resize(new_size, Image.Resampling.LANCZOS)
        
        # Gera nome único para a imagem
        image_filename = f"{Path(pdf_path).stem}_slide_{page_number}_{uuid.uuid4().hex[:8]}.png"
        image_path = os.path.join(self.temp_dir, image_filename)
        
        # Salva imagem otimizada
        pil_image.save(image_path, self.image_format, optimize=True)
        
        return image_path
    
    def _extract_table_text(self, table) -> str:
        """
        Extrai texto de uma tabela PowerPoint
        """
        if not table:
            return ""
            
        table_text = []
        
        try:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    cell_text = cell.text.strip() if cell.text else ""
                    row_text.append(cell_text)
                
                if any(cell for cell in row_text):  # Se há conteúdo na linha
                    table_text.append(" | ".join(row_text))
            
            return "\n".join(table_text)
        except Exception as e:
            logger.debug(f"Erro ao extrair texto da tabela: {e}")
            return ""
    
    def cleanup_temp_files(self, slide_images: List[str]):
        """
        Remove arquivos temporários de imagens
        """
        for image_path in slide_images:
            if image_path and os.path.exists(image_path):
                try:
                    os.unlink(image_path)
                except Exception as e:
                    logger.warning(f"Erro ao remover arquivo temporário {image_path}: {e}")
    
    async def convert_image_to_base64(self, image_path: str) -> Optional[str]:
        """
        Converte imagem para base64 para envio ao Gemini
        Otimizada para análise visual com tratamento robusto de formatos
        """
        if not image_path or not os.path.exists(image_path):
            return None
        
        try:
            with Image.open(image_path) as img:
                # Converte formatos problemáticos para RGB
                if img.mode in ('P', 'PA'):  # Palette com/sem transparência
                    if 'transparency' in img.info:
                        img = img.convert('RGBA')
                    else:
                        img = img.convert('RGB')
                elif img.mode not in ('RGB', 'RGBA'):
                    img = img.convert('RGB')
                
                # Verifica tamanho do arquivo
                original_size = os.path.getsize(image_path)
                
                # Se arquivo é muito grande, redimensiona mais agressivamente
                if original_size > settings.max_image_size_bytes:
                    # Calcula redução necessária
                    reduction_ratio = (settings.max_image_size_bytes / original_size) ** 0.5
                    new_size = tuple(int(dim * reduction_ratio) for dim in img.size)
                    img = img.resize(new_size, Image.Resampling.LANCZOS)
                
                # Limita dimensões máximas para análise eficiente
                max_dimension = min(self.max_image_dimension, 1536)  # Gemini funciona bem com até 1536px
                if max(img.size) > max_dimension:
                    ratio = max_dimension / max(img.size)
                    new_size = tuple(int(dim * ratio) for dim in img.size)
                    img = img.resize(new_size, Image.Resampling.LANCZOS)
                
                # Converte RGBA para RGB se necessário para JPEG
                if img.mode == 'RGBA':
                    # Cria fundo branco para imagens com transparência
                    rgb_img = Image.new('RGB', img.size, (255, 255, 255))
                    rgb_img.paste(img, mask=img.split()[-1])  # Usa canal alpha como máscara
                    img = rgb_img
                
                # Salva em buffer com qualidade otimizada
                buffer = io.BytesIO()
                
                # Ajusta qualidade baseada na configuração
                quality = 95 if settings.image_quality == "high" else 85 if settings.image_quality == "medium" else 75
                
                img.save(buffer, format='JPEG', quality=quality, optimize=True)
                buffer.seek(0)
                
                # Verifica tamanho final
                if buffer.getbuffer().nbytes > settings.max_image_size_bytes:
                    # Se ainda está grande, reduz qualidade
                    buffer = io.BytesIO()
                    img.save(buffer, format='JPEG', quality=60, optimize=True)
                    buffer.seek(0)
                
                # Converte para base64
                base64_string = base64.b64encode(buffer.getvalue()).decode()
                
                logger.debug(f"Imagem convertida: {len(base64_string)} caracteres base64")
                return base64_string
        
        except Exception as e:
            logger.error(f"Erro ao converter imagem para base64: {e}")
            return None
